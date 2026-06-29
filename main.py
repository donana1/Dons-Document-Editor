import os
import tkinter as tk
from tkinter import colorchooser, filedialog, font as tkfont, messagebox, ttk

import fitz
import PyPDF2
from docx import Document
from docx2pdf import convert as docx_to_pdf
from pdf2docx import Converter
from pdf2image import convert_from_path
from PIL import Image, ImageTk
from pptx import Presentation

# Professional Color Palette
COLOR_PRIMARY = "#1E3A8A"      # Deep Blue
COLOR_SECONDARY = "#3B82F6"    # Vibrant Blue
COLOR_BG_DARK = "#0F172A"      # Dark Slate for Sidebar
COLOR_BG_LIGHT = "#F8FAFC"     # Light Gray for Content Background
COLOR_TEXT_DARK = "#1E293B"    # Charcoal text
COLOR_TEXT_LIGHT = "#FFFFFF"   # White text
COLOR_BORDER = "#E2E8F0"       # Soft border gray


class FontPickerPopup(tk.Toplevel):
    def __init__(self, parent, fonts, on_select):
        super().__init__(parent)
        self.title("Select Font")
        self.geometry("340x400")
        self.resizable(False, False)
        self.configure(bg=COLOR_BG_LIGHT)
        self.transient(parent)
        self.grab_set()
        self.on_select = on_select
        self.fonts = fonts

        lbl = tk.Label(self, text="Choose System Font", font=("Segoe UI", 11, "bold"), bg=COLOR_BG_LIGHT, fg=COLOR_TEXT_DARK)
        lbl.pack(anchor="w", padx=12, pady=(12, 6))

        self.search_var = tk.StringVar()
        entry = tk.Entry(self, textvariable=self.search_var, font=("Segoe UI", 10), bg="white", fg=COLOR_TEXT_DARK, highlightbackground=COLOR_BORDER, highlightthickness=1, bd=0)
        entry.pack(fill="x", padx=12, pady=(0, 8), ipady=4)
        entry.bind("<KeyRelease>", self._filter)

        holder = tk.Frame(self, bg=COLOR_BG_LIGHT)
        holder.pack(fill="both", expand=True, padx=12, pady=(0, 12))

        self.listbox = tk.Listbox(holder, font=("Segoe UI", 10), bg="white", fg=COLOR_TEXT_DARK, bd=0, highlightbackground=COLOR_BORDER, highlightthickness=1, exportselection=False)
        self.listbox.pack(side="left", fill="both", expand=True)
        
        sb = ttk.Scrollbar(holder, orient="vertical", command=self.listbox.yview)
        sb.pack(side="right", fill="y")
        self.listbox.configure(yscrollcommand=sb.set)
        
        self.listbox.bind("<Double-Button-1>", self._choose)
        self.listbox.bind("<Return>", self._choose)
        self._set_items(self.fonts)

    def _set_items(self, items):
        self.listbox.delete(0, "end")
        for item in items:
            self.listbox.insert("end", item)
        if items:
            self.listbox.selection_set(0)

    def _filter(self, _event=None):
        q = self.search_var.get().strip().lower()
        if not q:
            self._set_items(self.fonts)
            return
        self._set_items([f for f in self.fonts if q in f.lower()])

    def _choose(self, _event=None):
        sel = self.listbox.curselection()
        if not sel:
            return
        self.on_select(self.listbox.get(sel[0]))
        self.destroy()


class DashboardApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Don's Document Editor")
        self.geometry("1400x850")
        self.minsize(1200, 750)
        self.configure(bg=COLOR_BG_LIGHT)

        # Style configurations
        self.style = ttk.Style()
        self.style.theme_use("clam")
        self.style.configure("TFrame", background=COLOR_BG_LIGHT)
        self.style.configure("Sidebar.TFrame", background=COLOR_BG_DARK)
        self.style.configure("Card.TFrame", background="white", relief="solid", borderwidth=1)
        
        # Application States
        self.merge_files = []
        
        # Editor Workspace States
        self.doc = None
        self.page_idx = 0
        self.scale = 1.25
        self.page_img = None
        self.page_img_tk = None
        self.thumb_refs = []
        self.thumb_frames = []

        self.mode = tk.StringVar(value="pointer")
        self.pen_width = tk.IntVar(value=4)
        self.pen_color = "#3B82F6"
        self.pending_add_text = False

        self.strokes = []
        self.text_boxes = []
        self.inline_edits = []  
        self.active_stroke = None
        self.active_text_idx = None
        self.dragging_text_idx = None
        self.resizing_text_idx = None
        self.drag_last = (0, 0)
        self.edit_widget = None

        self.system_fonts = sorted(tkfont.families())
        self.font_family = tk.StringVar(value="Arial")
        self.font_size = tk.IntVar(value=14)
        self.bold = tk.BooleanVar(value=False)
        self.italic = tk.BooleanVar(value=False)
        self.strike = tk.BooleanVar(value=False)

        self._build_main_structure()

    def _build_main_structure(self):
        # 1. Left Sidebar
        self.sidebar = tk.Frame(self, bg=COLOR_BG_DARK, width=260)
        self.sidebar.pack(side="left", fill="y")
        self.sidebar.pack_propagate(False)

        brand_lbl = tk.Label(self.sidebar, text="Don's Document Studio", font=("Segoe UI", 13, "bold"), fg=COLOR_TEXT_LIGHT, bg=COLOR_BG_DARK)
        brand_lbl.pack(anchor="w", padx=20, pady=(24, 4))
        sub_brand = tk.Label(self.sidebar, text="Professional Workspace", font=("Segoe UI", 9), fg="#94A3B8", bg=COLOR_BG_DARK)
        sub_brand.pack(anchor="w", padx=20, pady=(0, 30))

        # Sidebar Buttons
        self.nav_buttons = {}
        routes = [("converter", "Format Converter"), ("editor", "Interactive Editor"), ("merger", "Merge Documents")]
        for route_id, label in routes:
            btn = tk.Button(
                self.sidebar, text=label, font=("Segoe UI", 11),
                fg=COLOR_TEXT_LIGHT, bg=COLOR_BG_DARK, activebackground=COLOR_PRIMARY, activeforeground=COLOR_TEXT_LIGHT,
                bd=0, anchor="w", padx=20, pady=12, cursor="hand2",
                command=lambda r=route_id: self.switch_view(r)
            )
            btn.pack(fill="x", padx=10, pady=2)
            self.nav_buttons[route_id] = btn

        footer_lbl = tk.Label(self.sidebar, text="Created by Don", font=("Segoe UI", 9), fg="#64748B", bg=COLOR_BG_DARK)
        footer_lbl.pack(side="bottom", anchor="w", padx=20, pady=16)

        # 2. Main content container
        self.content_viewport = tk.Frame(self, bg=COLOR_BG_LIGHT)
        self.content_viewport.pack(side="right", fill="both", expand=True)

        # Initialize internal views
        self.views = {}
        self._init_converter_view()
        self._init_editor_view()
        self._init_merger_view()

        # Default View
        self.switch_view("converter")

    def switch_view(self, route_id):
        for v in self.views.values():
            v.pack_forget()
        
        for r_id, btn in self.nav_buttons.items():
            if r_id == route_id:
                btn.configure(bg=COLOR_PRIMARY, font=("Segoe UI", 11, "bold"))
            else:
                btn.configure(bg=COLOR_BG_DARK, font=("Segoe UI", 11, "normal"))

        self.views[route_id].pack(fill="both", expand=True)

    def _create_header(self, parent, title, subtitle):
        header_frame = tk.Frame(parent, bg=COLOR_BG_LIGHT)
        header_frame.pack(fill="x", padx=30, pady=(24, 16))
        tk.Label(header_frame, text=title, font=("Segoe UI", 20, "bold"), fg=COLOR_TEXT_DARK, bg=COLOR_BG_LIGHT).pack(anchor="w")
        tk.Label(header_frame, text=subtitle, font=("Segoe UI", 10), fg="#64748B", bg=COLOR_BG_LIGHT).pack(anchor="w", pady=(2, 0))
        return header_frame

    def _init_converter_view(self):
        v = tk.Frame(self.content_viewport, bg=COLOR_BG_LIGHT)
        self.views["converter"] = v
        self._create_header(v, "Format Converter", "Transform documents between common corporate definitions quickly.")

        card = tk.Frame(v, bg="white", highlightbackground=COLOR_BORDER, highlightthickness=1)
        card.pack(fill="both", expand=True, padx=30, pady=(0, 30))

        grid = tk.Frame(card, bg="white")
        grid.place(relx=0.5, rely=0.45, anchor="center")

        actions = [
            ("Word to PDF", self.word_to_pdf),
            ("PDF to Word", self.pdf_to_word),
            ("PowerPoint to PDF", self.ppt_to_pdf),
            ("PDF to PowerPoint", self.pdf_to_ppt),
            ("JPEG to PDF", lambda: self.image_to_pdf(("*.jpg *.jpeg", "*.jpg", "*.jpeg"))),
            ("PDF to JPEG", lambda: self.pdf_to_images("JPEG")),
            ("PNG to PDF", lambda: self.image_to_pdf(("*.png",))),
            ("PDF to PNG", lambda: self.pdf_to_images("PNG")),
        ]
        
        for i, (label, cmd) in enumerate(actions):
            r, c = divmod(i, 2)
            btn = tk.Button(
                grid, text=label, font=("Segoe UI", 11), 
                bg=COLOR_PRIMARY, fg=COLOR_TEXT_LIGHT,
                activebackground=COLOR_SECONDARY, activeforeground=COLOR_TEXT_LIGHT, bd=0,
                width=24, height=2, cursor="hand2", command=cmd
            )
            btn.grid(row=r, column=c, padx=16, pady=16)

    def _init_merger_view(self):
        v = tk.Frame(self.content_viewport, bg=COLOR_BG_LIGHT)
        self.views["merger"] = v
        self._create_header(v, "Merge Documents", "Combine individual PDF elements cleanly into a single structured file.")

        card = tk.Frame(v, bg="white", highlightbackground=COLOR_BORDER, highlightthickness=1)
        card.pack(fill="both", expand=True, padx=30, pady=(0, 30))

        left_pane = tk.Frame(card, bg="white")
        left_pane.pack(side="left", fill="both", expand=True, padx=20, pady=20)
        
        tk.Label(left_pane, text="Document Queue", font=("Segoe UI", 11, "bold"), fg=COLOR_TEXT_DARK, bg="white").pack(anchor="w", pady=(0, 8))
        
        self.merge_list = tk.Listbox(left_pane, font=("Segoe UI", 11), bg=COLOR_BG_LIGHT, fg=COLOR_TEXT_DARK, bd=0, highlightbackground=COLOR_BORDER, highlightthickness=1)
        self.merge_list.pack(fill="both", expand=True)

        right_pane = tk.Frame(card, bg="white", width=220)
        right_pane.pack(side="right", fill="y", padx=20, pady=20)
        right_pane.pack_propagate(False)

        btns = [
            ("Add PDFs", self.add_merge_files, COLOR_PRIMARY),
            ("Remove Selected", self.remove_merge_file, "#EF4444"),
            ("Move Up", lambda: self.move_merge_item(-1), COLOR_SECONDARY),
            ("Move Down", lambda: self.move_merge_item(1), COLOR_SECONDARY),
            ("Execute Merge", self.merge_pdfs, "#10B981")
        ]
        for label, cmd, col in btns:
            b = tk.Button(
                right_pane, text=label, font=("Segoe UI", 10, "bold" if "Merge" in label else "normal"), 
                bg=col, fg=COLOR_TEXT_LIGHT, activebackground=COLOR_BG_DARK, activeforeground=COLOR_TEXT_LIGHT,
                bd=0, height=2, cursor="hand2", command=cmd
            )
            b.pack(fill="x", pady=5)

    def _init_editor_view(self):
        v = tk.Frame(self.content_viewport, bg=COLOR_BG_LIGHT)
        self.views["editor"] = v
        
        # Tool Controls Context Top-bar
        bar = tk.Frame(v, bg="white", highlightbackground=COLOR_BORDER, highlightthickness=1)
        bar.pack(fill="x", padx=20, pady=(15, 10))

        # Document loading buttons
        tk.Button(bar, text="Open PDF", font=("Segoe UI", 10, "bold"), bg=COLOR_PRIMARY, fg=COLOR_TEXT_LIGHT, bd=0, padx=12, pady=6, command=self.open_pdf, cursor="hand2").pack(side="left", padx=6, pady=6)
        tk.Button(bar, text="Save Changes", font=("Segoe UI", 10, "bold"), bg="#10B981", fg=COLOR_TEXT_LIGHT, bd=0, padx=12, pady=6, command=self.save_pdf, cursor="hand2").pack(side="left", padx=4, pady=6)
        
        tk.Frame(bar, bg=COLOR_BORDER, width=1).pack(side="left", fill="y", padx=8, pady=6)

        # Mode Toggles (With Text Modifier Included)
        tk.Radiobutton(bar, text="Pointer Tool", font=("Segoe UI", 10), value="pointer", variable=self.mode, bg="white", selectcolor="white", fg=COLOR_TEXT_DARK).pack(side="left", padx=4)
        tk.Radiobutton(bar, text="Text Modifier", font=("Segoe UI", 10), value="modifier", variable=self.mode, bg="white", selectcolor="white", fg=COLOR_TEXT_DARK).pack(side="left", padx=4)
        tk.Radiobutton(bar, text="Draw Pen", font=("Segoe UI", 10), value="pen", variable=self.mode, bg="white", selectcolor="white", fg=COLOR_TEXT_DARK).pack(side="left", padx=4)
        tk.Radiobutton(bar, text="Eraser", font=("Segoe UI", 10), value="eraser", variable=self.mode, bg="white", selectcolor="white", fg=COLOR_TEXT_DARK).pack(side="left", padx=4)

        tk.Frame(bar, bg=COLOR_BORDER, width=1).pack(side="left", fill="y", padx=8, pady=6)
        
        tk.Button(bar, text="＋ Text Box", font=("Segoe UI", 10), bg="white", fg=COLOR_PRIMARY, highlightbackground=COLOR_PRIMARY, highlightthickness=1, bd=0, padx=10, command=self.queue_add_textbox, cursor="hand2").pack(side="left", padx=4)
        tk.Button(bar, text="Color Picker", font=("Segoe UI", 10), bg="white", fg=COLOR_TEXT_DARK, highlightbackground=COLOR_BORDER, highlightthickness=1, bd=0, padx=10, command=self.pick_pen_color, cursor="hand2").pack(side="left", padx=4)
        
        tk.Label(bar, text="Width", font=("Segoe UI", 9), fg="#64748B", bg="white").pack(side="left", padx=(6, 2))
        tk.Scale(bar, from_=1, to=40, variable=self.pen_width, orient="horizontal", length=80, bg="white", bd=0, highlightthickness=0).pack(side="left", padx=2)

        tk.Frame(bar, bg=COLOR_BORDER, width=1).pack(side="left", fill="y", padx=8, pady=6)

        # Typography contextual sub-bar
        tk.Button(bar, text="Font Family", font=("Segoe UI", 10), bg="white", fg=COLOR_TEXT_DARK, bd=0, command=self.open_font_popup, cursor="hand2").pack(side="left", padx=2)
        tk.Label(bar, textvariable=self.font_family, font=("Segoe UI", 10, "italic"), bg=COLOR_BG_LIGHT, fg=COLOR_PRIMARY, width=12).pack(side="left", padx=2)
        
        size_box = ttk.Spinbox(bar, from_=6, to=120, width=4, textvariable=self.font_size, command=self.apply_style_to_active_text)
        size_box.pack(side="left", padx=2)
        size_box.bind("<Return>", self.apply_style_to_active_text)
        
        tk.Checkbutton(bar, text="B", font=("Segoe UI", 10, "bold"), variable=self.bold, command=self.apply_style_to_active_text, bg="white", selectcolor="white").pack(side="left", padx=1)
        tk.Checkbutton(bar, text="I", font=("Segoe UI", 10, "italic"), variable=self.italic, command=self.apply_style_to_active_text, bg="white", selectcolor="white").pack(side="left", padx=1)
        tk.Checkbutton(bar, text="S", font=("Segoe UI", 10, "overstrike"), variable=self.strike, command=self.apply_style_to_active_text, bg="white", selectcolor="white").pack(side="left", padx=1)

        # Layout Main Pane splits
        work_area = tk.Frame(v, bg=COLOR_BG_LIGHT)
        work_area.pack(fill="both", expand=True, padx=20, pady=(0, 20))
        
        # Left Side Canvas Core
        self.left_canvas_frame = tk.Frame(work_area, bg="white", highlightbackground=COLOR_BORDER, highlightthickness=1)
        self.left_canvas_frame.pack(side="left", fill="both", expand=True)

        self.canvas = tk.Canvas(self.left_canvas_frame, bg="white", highlightthickness=0)
        self.canvas.pack(fill="both", expand=True)
        self.canvas.bind("<Button-1>", self.on_canvas_down)
        self.canvas.bind("<B1-Motion>", self.on_canvas_drag)
        self.canvas.bind("<ButtonRelease-1>", self.on_canvas_up)
        self.canvas.bind("<Double-Button-1>", self.on_canvas_double_click)

        # Right Side Thumbnails panel Structure
        self.right_panel = tk.Frame(work_area, bg="white", width=260, highlightbackground=COLOR_BORDER, highlightthickness=1)
        self.right_panel.pack(side="right", fill="y", padx=(15, 0))
        self.right_panel.pack_propagate(False)

        self.thumb_canvas = tk.Canvas(self.right_panel, bg=COLOR_BG_LIGHT, highlightthickness=0)
        self.thumb_sb = ttk.Scrollbar(self.right_panel, orient="vertical", command=self.thumb_canvas.yview)
        self.thumb_canvas.configure(yscrollcommand=self.thumb_sb.set)
        
        self.thumb_sb.pack(side="right", fill="y")
        self.thumb_canvas.pack(side="top", fill="both", expand=True)
        
        self.thumb_inner = tk.Frame(self.thumb_canvas, bg=COLOR_BG_LIGHT)
        self.thumb_canvas.create_window((0, 0), window=self.thumb_inner, anchor="nw")
        self.thumb_inner.bind("<Configure>", lambda _e: self.thumb_canvas.configure(scrollregion=self.thumb_canvas.bbox("all")))

        # Rearranging Operations Deck Bottom-right
        deck = tk.Frame(self.right_panel, bg="white", highlightbackground=COLOR_BORDER, highlightthickness=1)
        deck.pack(side="bottom", fill="x", pady=0)
        
        tk.Button(deck, text="▲ Move Page Up", font=("Segoe UI", 10), bg=COLOR_SECONDARY, fg=COLOR_TEXT_LIGHT, bd=0, height=2, command=lambda: self.move_active_page_item(-1), cursor="hand2").pack(side="left", fill="x", expand=True, padx=4, pady=6)
        tk.Button(deck, text="▼ Move Page Down", font=("Segoe UI", 10), bg=COLOR_SECONDARY, fg=COLOR_TEXT_LIGHT, bd=0, height=2, command=lambda: self.move_active_page_item(1), cursor="hand2").pack(side="left", fill="x", expand=True, padx=4, pady=6)

    def move_active_page_item(self, step):
        if not self.doc or len(self.doc) <= 1:
            return
        
        src_idx = self.page_idx
        dest_idx = src_idx + step
        if dest_idx < 0 or dest_idx >= len(self.doc):
            return

        self.doc.move_page(src_idx, dest_idx)

        moving_stroke = self.strokes.pop(src_idx)
        self.strokes.insert(dest_idx, moving_stroke)

        moving_boxes = self.text_boxes.pop(src_idx)
        self.text_boxes.insert(dest_idx, moving_boxes)

        if hasattr(self, 'inline_edits') and len(self.inline_edits) > src_idx:
            moving_inlines = self.inline_edits.pop(src_idx)
            self.inline_edits.insert(dest_idx, moving_inlines)

        self.page_idx = dest_idx

        self.render_page()
        self.render_thumbnails()

        self.thumb_canvas.update_idletasks()
        if dest_idx < len(self.thumb_frames):
            target_frame = self.thumb_frames[dest_idx]
            y_pos = target_frame.winfo_y() / max(1, self.thumb_inner.winfo_height())
            self.thumb_canvas.yview_moveto(max(0.0, y_pos - 0.2))

    def open_font_popup(self):
        FontPickerPopup(self, self.system_fonts, self.set_font_family)

    def set_font_family(self, family):
        self.font_family.set(family)
        self.apply_style_to_active_text()

    def queue_add_textbox(self):
        self.pending_add_text = True
        self.mode.set("pointer")

    def pick_pen_color(self):
        c = colorchooser.askcolor(color=self.pen_color)[1]
        if c:
            self.pen_color = c

    def open_pdf(self):
        path = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if not path:
            return
        self.doc = fitz.open(path)
        self.page_idx = 0
        self.strokes = [[] for _ in range(len(self.doc))]
        self.text_boxes = [[] for _ in range(len(self.doc))]
        self.inline_edits = [[] for _ in range(len(self.doc))]
        self.render_page()
        self.render_thumbnails()

    def render_page(self):
        self.canvas.delete("all")
        if not self.doc:
            return
        page = self.doc[self.page_idx]
        pix = page.get_pixmap(matrix=fitz.Matrix(self.scale, self.scale))
        self.page_img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        self.page_img_tk = ImageTk.PhotoImage(self.page_img)
        self.canvas.create_image(0, 0, anchor="nw", image=self.page_img_tk)
        self.canvas.config(scrollregion=(0, 0, pix.width, pix.height))
        self._draw_overlays()

    def _draw_overlays(self):
        if hasattr(self, 'inline_edits') and len(self.inline_edits) > self.page_idx:
            for edit in self.inline_edits[self.page_idx]:
                rx0, ry0, rx1, ry1 = [v * self.scale for v in edit["bbox"]]
                self.canvas.create_rectangle(rx0, ry0, rx1, ry1, fill="white", outline="white")
                self.canvas.create_text(
                    rx0, ry0, anchor="nw", text=edit["text"],
                    font=(edit.get("font", "Arial"), int(edit.get("size", 11) * self.scale)), fill="black"
                )

        for s in self.strokes[self.page_idx]:
            pts = s["points"]
            if len(pts) < 2:
                continue
            flat = [v for p in pts for v in p]
            self.canvas.create_line(*flat, fill=s["color"], width=s["width"], capstyle="round", smooth=True)

        for t in self.text_boxes[self.page_idx]:
            x, y, w, h = t["x"], t["y"], t["w"], t["h"]
            self.canvas.create_rectangle(x, y, x + w, y + h, outline=COLOR_SECONDARY, width=1)
            style = []
            if t.get("bold"): style.append("bold")
            if t.get("italic"): style.append("italic")
            font_tuple = (t.get("font", "Arial"), t.get("size", 14), " ".join(style) if style else "normal")
            text_id = self.canvas.create_text(x + 4, y + 4, anchor="nw", text=t.get("text", ""), width=max(10, w - 8), font=font_tuple, fill="black")
            if t.get("strike"):
                bb = self.canvas.bbox(text_id)
                if bb:
                    self.canvas.create_line(bb[0], (bb[1] + bb[3]) / 2, bb[2], (bb[1] + bb[3]) / 2, fill="black", width=1)
            self.canvas.create_rectangle(x + w - 8, y + h - 8, x + w, y + h, fill=COLOR_SECONDARY, outline=COLOR_SECONDARY)

    def _hit_text(self, x, y):
        boxes = self.text_boxes[self.page_idx]
        for i in range(len(boxes) - 1, -1, -1):
            b = boxes[i]
            if b["x"] <= x <= b["x"] + b["w"] and b["y"] <= y <= b["y"] + b["h"]:
                return i
        return None

    @staticmethod
    def _near_handle(b, x, y):
        return b["x"] + b["w"] - 12 <= x <= b["x"] + b["w"] + 2 and b["y"] + b["h"] - 12 <= y <= b["y"] + b["h"] + 2

    def on_canvas_down(self, event):
        if not self.doc:
            return
        x, y = event.x, event.y
        idx = self.page_idx

        if self.mode.get() == "modifier":
            page = self.doc[idx]
            native_x, native_y = x / self.scale, y / self.scale
            text_page = page.get_text("words") 
            matched_word = None
            for w in text_page:
                x0, y0, x1, y1, word_str, block_no, line_no, word_no = w
                if x0 <= native_x <= x1 and y0 <= native_y <= y1:
                    matched_word = w
                    break
            
            if matched_word:
                x0, y0, x1, y1, word_str, _, _, _ = matched_word
                if self.edit_widget: self.edit_widget.destroy()
                
                rx0, ry0 = x0 * self.scale, y0 * self.scale
                rw, rh = (x1 - x0) * self.scale, (y1 - y0) * self.scale
                
                txt = tk.Text(self.canvas, bd=0, wrap="none", font=("Arial", max(9, int(rh * 0.75))), highlightbackground=COLOR_PRIMARY, highlightthickness=1)
                txt.insert("1.0", word_str)
                txt.place(x=rx0, y=ry0, width=max(60, rw + 20), height=rh + 4)
                txt.focus_set()
                self.edit_widget = txt
                
                def save_inline_modification(_e=None):
                    new_val = txt.get("1.0", "end-1c").strip()
                    txt.destroy()
                    self.edit_widget = None
                    if new_val != word_str:
                        self.inline_edits[idx].append({
                            "bbox": (x0, y0, x1, y1),
                            "text": new_val,
                            "font": "Arial",
                            "size": (y1 - y0) * 0.85
                        })
                    self.render_page()

                txt.bind("<FocusOut>", save_inline_modification)
                txt.bind("<Escape>", lambda _e: [txt.destroy(), setattr(self, 'edit_widget', None)])
                txt.bind("<Return>", save_inline_modification)
            return

        if self.pending_add_text:
            self.pending_add_text = False
            self.text_boxes[idx].append({
                "x": x, "y": y, "w": 200, "h": 80, "text": "Click to edit text",
                "font": self.font_family.get(), "size": int(self.font_size.get()),
                "bold": bool(self.bold.get()), "italic": bool(self.italic.get()), "strike": bool(self.strike.get()),
            })
            self.active_text_idx = len(self.text_boxes[idx]) - 1
            self.render_page()
            return

        if self.mode.get() == "pen":
            self.active_stroke = {"points": [(x, y)], "color": self.pen_color, "width": int(self.pen_width.get())}
            self.strokes[idx].append(self.active_stroke)
            return
        if self.mode.get() == "eraser":
            self.erase_strokes(x, y)
            self.render_page()
            return

        hit = self._hit_text(x, y)
        self.dragging_text_idx = None
        self.resizing_text_idx = None
        self.active_text_idx = hit
        if hit is None:
            return
        b = self.text_boxes[idx][hit]
        self.drag_last = (x, y)
        if self._near_handle(b, x, y):
            self.resizing_text_idx = hit
        else:
            self.dragging_text_idx = hit
        self.sync_toolbar_from_box(b)

    def on_canvas_drag(self, event):
        if not self.doc:
            return
        x, y = event.x, event.y
        idx = self.page_idx
        mode = self.mode.get()

        if mode == "pen" and self.active_stroke is not None:
            self.active_stroke["points"].append((x, y))
            self.render_page()
            return
        if mode == "eraser":
            self.erase_strokes(x, y)
            self.render_page()
            return

        if self.dragging_text_idx is not None:
            b = self.text_boxes[idx][self.dragging_text_idx]
            dx, dy = x - self.drag_last[0], y - self.drag_last[1]
            b["x"] += dx
            b["y"] += dy
            self.drag_last = (x, y)
            self.render_page()
        elif self.resizing_text_idx is not None:
            b = self.text_boxes[idx][self.resizing_text_idx]
            b["w"] = max(60, x - b["x"])
            b["h"] = max(30, y - b["y"])
            self.render_page()

    def on_canvas_up(self, _event):
        self.active_stroke = None
        self.dragging_text_idx = None
        self.resizing_text_idx = None

    def on_canvas_double_click(self, event):
        if self.mode.get() != "pointer" or not self.doc:
            return
        hit = self._hit_text(event.x, event.y)
        if hit is None:
            return
        self.edit_textbox(hit)

    def edit_textbox(self, idx):
        if self.edit_widget is not None:
            self.edit_widget.destroy()
        b = self.text_boxes[self.page_idx][idx]
        
        txt = tk.Text(self.canvas, bd=0, wrap="word", font=(b.get("font", "Arial"), 11), highlightbackground=COLOR_SECONDARY, highlightthickness=1)
        txt.insert("1.0", b.get("text", ""))
        txt.place(x=b["x"], y=b["y"], width=b["w"], height=b["h"])
        txt.focus_set()
        self.edit_widget = txt

        def finish(_event=None):
            b["text"] = txt.get("1.0", "end-1c")
            txt.destroy()
            self.edit_widget = None
            self.render_page()

        txt.bind("<FocusOut>", finish)
        txt.bind("<Escape>", finish)

    def sync_toolbar_from_box(self, b):
        self.font_family.set(b.get("font", "Arial"))
        self.font_size.set(int(b.get("size", 14)))
        self.bold.set(bool(b.get("bold", False)))
        self.italic.set(bool(b.get("italic", False)))
        self.strike.set(bool(b.get("strike", False)))

    def apply_style_to_active_text(self, _event=None):
        if not self.doc or self.active_text_idx is None:
            return
        b = self.text_boxes[self.page_idx][self.active_text_idx]
        b["font"] = self.font_family.get()
        b["size"] = int(self.font_size.get())
        b["bold"] = bool(self.bold.get())
        b["italic"] = bool(self.italic.get())
        b["strike"] = bool(self.strike.get())
        self.render_page()

    def erase_strokes(self, x, y):
        r = max(8, int(self.pen_width.get()) // 2 + 5)
        self.strokes[self.page_idx] = [
            s for s in self.strokes[self.page_idx]
            if not any((px - x) ** 2 + (py - y) ** 2 <= r * r for px, py in s["points"])
        ]

    def render_thumbnails(self):
        for w in self.thumb_inner.winfo_children():
            w.destroy()
    
        self.thumb_refs = []
        self.thumb_frames = []
        if not self.doc:
            return
        
        for i in range(len(self.doc)):
            is_active = (i == self.page_idx)
            row = tk.Frame(self.thumb_inner, bg="white", highlightbackground=COLOR_PRIMARY if is_active else COLOR_BORDER, highlightthickness=2 if is_active else 1)
            row.pack(fill="x", padx=10, pady=6)
            self.thumb_frames.append(row)
            
            lbl_num = tk.Label(row, text=f"Page {i+1}", font=("Segoe UI", 9, "bold" if is_active else "normal"), fg=COLOR_PRIMARY if is_active else COLOR_TEXT_DARK, bg="white")
            lbl_num.pack(anchor="w", padx=6, pady=2)

            pix = self.doc[i].get_pixmap(matrix=fitz.Matrix(0.22, 0.22))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            tk_img = ImageTk.PhotoImage(img)
            self.thumb_refs.append(tk_img)
            
            lbl = tk.Label(row, image=tk_img, bd=0, bg="white")
            lbl.pack(padx=6, pady=(0, 6))

            rot = tk.Button(row, text="↻", font=("Segoe UI", 8), bg=COLOR_BG_LIGHT, fg=COLOR_TEXT_DARK, bd=1, command=lambda idx=i: self.rotate_page(idx), cursor="hand2")
            row.bind("<Enter>", lambda _e, b=rot: b.place(relx=1.0, rely=0.0, anchor="ne", x=-4, y=4))
            row.bind("<Leave>", lambda _e, b=rot: b.place_forget())

            for widget in (row, lbl, lbl_num):
                widget.bind("<ButtonPress-1>", lambda _e, idx=i: self.thumb_press(idx))

    def thumb_press(self, idx):
        self.page_idx = idx
        self.render_page()
        self.render_thumbnails()

    def rotate_page(self, idx):
        if not self.doc:
            return
        p = self.doc[idx]
        p.set_rotation((p.rotation + 90) % 360)
        self.render_page()
        self.render_thumbnails()

    def save_pdf(self):
        if not self.doc:
            return
        if self.edit_widget is not None:
            self.edit_widget.event_generate("<FocusOut>")
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
        if not out:
            return
        for i in range(len(self.doc)):
            page = self.doc[i]
            
            if hasattr(self, 'inline_edits') and len(self.inline_edits) > i:
                for edit in self.inline_edits[i]:
                    x0, y0, x1, y1 = edit["bbox"]
                    page.draw_rect(fitz.Rect(x0, y0, x1, y1), color=(1, 1, 1), fill=(1, 1, 1))
                    page.insert_text(fitz.Point(x0, y1 - 2), edit["text"], fontsize=edit["size"], color=(0, 0, 0))

            for s in self.strokes[i]:
                rgb = tuple(int(s["color"].lstrip("#")[j : j + 2], 16) / 255 for j in (0, 2, 4))
                w = max(0.2, s["width"] / self.scale)
                for a, b in zip(s["points"], s["points"][1:]):
                    p1 = fitz.Point(a[0] / self.scale, a[1] / self.scale)
                    p2 = fitz.Point(b[0] / self.scale, b[1] / self.scale)
                    page.draw_line(p1, p2, color=rgb, width=w)
            for t in self.text_boxes[i]:
                x, y = t["x"] / self.scale, t["y"] / self.scale
                size = float(t.get("size", 14))
                for ln, line in enumerate(t.get("text", "").splitlines()):
                    y0 = y + ln * size * 1.2
                    page.insert_text(fitz.Point(x, y0), line, fontsize=size, color=(0, 0, 0))
                    if t.get("strike") and line:
                        lw = max(1.0, len(line) * size * 0.5)
                        ys = y0 - size * 0.3
                        page.draw_line(fitz.Point(x, ys), fitz.Point(x + lw, ys), color=(0, 0, 0), width=max(0.2, size * 0.06))
        self.doc.save(out)
        messagebox.showinfo("Saved", "Edited document output written successfully.")

    def word_to_pdf(self):
        src = filedialog.askopenfilename(filetypes=[("Word", "*.docx")])
        if not src: return
        try:
            docx_to_pdf(src)
            messagebox.showinfo("Success", "Word to PDF conversion complete.")
        except Exception as e: messagebox.showerror("Error", str(e))

    def pdf_to_word(self):
        src = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if not src: return
        out = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word", "*.docx")])
        if not out: return
        try:
            cv = Converter(src)
            cv.convert(out, start=0, end=None)
            cv.close()
            messagebox.showinfo("Success", "PDF to Word conversion complete.")
        except Exception as e: messagebox.showerror("Error", str(e))

    def ppt_to_pdf(self):
        src = filedialog.askopenfilename(filetypes=[("PowerPoint", "*.pptx")])
        if not src: return
        try:
            base, _ = os.path.splitext(src)
            prs = Presentation(src)
            tmp_docx = base + "_ppt_content.docx"
            doc = Document()
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        doc.add_paragraph(shape.text)
            doc.save(tmp_docx)
            docx_to_pdf(tmp_docx)
            messagebox.showinfo("Success", "PowerPoint to PDF conversion complete.")
        except Exception as e: messagebox.showerror("Error", str(e))

    def pdf_to_ppt(self):
        src = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if not src: return
        out = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")])
        if not out: return
        try:
            images = convert_from_path(src)
            prs = Presentation()
            base, _ = os.path.splitext(out)
            for n, img in enumerate(images):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                temp_img = f"{base}_tmp_{n}.png"
                img.save(temp_img, "PNG")
                slide.shapes.add_picture(temp_img, 0, 0, width=prs.slide_width, height=prs.slide_height)
                os.remove(temp_img)
            prs.save(out)
            messagebox.showinfo("Success", "PDF to PowerPoint conversion complete.")
        except Exception as e: messagebox.showerror("Error", str(e))

    def image_to_pdf(self, pattern_tuple):
        src = filedialog.askopenfilename(filetypes=[("Image", " ".join(pattern_tuple))])
        if not src: return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if not out: return
        try:
            img = Image.open(src).convert("RGB")
            img.save(out)
            messagebox.showinfo("Success", "Image to PDF conversion complete.")
        except Exception as e: messagebox.showerror("Error", str(e))

    def pdf_to_images(self, mode):
        src = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if not src: return
        out_dir = filedialog.askdirectory(title="Select output folder")
        if not out_dir: return
        try:
            images = convert_from_path(src)
            ext = "jpg" if mode == "JPEG" else "png"
            base = os.path.splitext(os.path.basename(src))[0]
            for i, img in enumerate(images, start=1):
                out = os.path.join(out_dir, f"{base}_page_{i}.{ext}")
                img.save(out, mode)
            messagebox.showinfo("Success", f"PDF to {mode} conversion complete.")
        except Exception as e: messagebox.showerror("Error", str(e))

    def add_merge_files(self):
        paths = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        if not paths: return
        for p in paths:
            self.merge_files.append(p)
            self.merge_list.insert("end", f" 📄  {os.path.basename(p)}")

    def remove_merge_file(self):
        sel = self.merge_list.curselection()
        if not sel: return
        i = sel[0]
        self.merge_list.delete(i)
        self.merge_files.pop(i)

    def move_merge_item(self, step):
        sel = self.merge_list.curselection()
        if not sel: return
        i = sel[0]
        j = i + step
        if j < 0 or j >= len(self.merge_files): return
        self.merge_files[i], self.merge_files[j] = self.merge_files[j], self.merge_files[i]
        self.merge_list.delete(0, "end")
        for p in self.merge_files:
            self.merge_list.insert("end", f" 📄  {os.path.basename(p)}")
        self.merge_list.selection_set(j)

    def merge_pdfs(self):
        if not self.merge_files:
            messagebox.showwarning("Warning", "The file queue is completely empty.")
            return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
        if not out: return
        try:
            merger = PyPDF2.PdfMerger()
            for p in self.merge_files:
                merger.append(p)
            merger.write(out)
            merger.close()
            messagebox.showinfo("Success", "Documents merged successfully.")
        except Exception as e: messagebox.showerror("Error", str(e))


if __name__ == "__main__":
    app = DashboardApp()
    app.mainloop()
